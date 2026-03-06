"""Generate PBRToonBase technical sharing PPT - Chinese, minimal code, B&W theme."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Black/White/Gray palette ===
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)      # light gray card
BG_MID = RGBColor(0xE8, 0xE8, 0xE8)         # mid gray
BLACK = RGBColor(0x1A, 0x1A, 0x1A)           # near black
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)       # dark gray text
MID_GRAY = RGBColor(0x66, 0x66, 0x66)        # body text
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)      # subtle text
ACCENT = RGBColor(0x33, 0x33, 0x33)          # accent = dark gray
BORDER = RGBColor(0xCC, 0xCC, 0xCC)          # border color

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BG_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(0), rounded=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(st, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height, text, sz=18, color=DARK_GRAY,
       bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_para(tf, text, sz=16, color=DARK_GRAY, bold=False, font="Microsoft YaHei",
             align=PP_ALIGN.LEFT, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_before = sp
    return p


def flow_box(slide, left, top, width, height, text, fill=BG_LIGHT, border=BORDER, sz=13):
    shape = add_rect(slide, left, top, width, height, fill, border, Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = DARK_GRAY
        p.font.bold = (i == 0)
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
    return shape


def arrow_down(slide, cx, top, length):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - Inches(0.1), top, Inches(0.2), length)
    s.fill.solid()
    s.fill.fore_color.rgb = LIGHT_GRAY
    s.line.fill.background()


def arrow_right(slide, left, cy, length):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, cy - Inches(0.08), length, Inches(0.16))
    s.fill.solid()
    s.fill.fore_color.rgb = LIGHT_GRAY
    s.line.fill.background()


def num_circle(slide, left, top, num):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.4), Inches(0.4))
    s.fill.solid()
    s.fill.fore_color.rgb = BLACK
    s.line.fill.background()
    p = s.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(15)
    p.font.color.rgb = BG_WHITE
    p.font.bold = True
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def title_bar(slide, title, subtitle=""):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()
    tb(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7), title, sz=30, color=BLACK, bold=True)
    if subtitle:
        tb(slide, Inches(0.8), Inches(0.72), Inches(11), Inches(0.35), subtitle, sz=15, color=LIGHT_GRAY)


def bullet_card(slide, x, y, w, h, title, items):
    """Card with title + bullet items."""
    add_rect(slide, x, y, w, h, BG_LIGHT, BORDER, Pt(1))
    tb(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.35),
       title, sz=16, color=BLACK, bold=True)
    # divider
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.48),
                               w - Inches(0.4), Inches(0.015))
    d.fill.solid()
    d.fill.fore_color.rgb = BG_MID
    d.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.52), w - Inches(0.5), h - Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = MID_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(5)


def formula_box(slide, left, top, width, height, text, sz=13):
    """Minimal inline formula/code, gray bg."""
    shape = add_rect(slide, left, top, width, height, BG_LIGHT, BORDER, Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    for i, line in enumerate(text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Consolas"
        p.space_before = Pt(2)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Top accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLACK
bar.line.fill.background()

tb(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
   "PBRToonBase", sz=56, color=BLACK, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
tb(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.7),
   "PBR + Toon 混合渲染管线技术分享", sz=26, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(2.333), Inches(0.025))
div.fill.solid()
div.fill.fore_color.rgb = BLACK
div.line.fill.background()

tb(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5),
   "Arknights: Endfield  |  Goo Engine 4.4  |  Blender 4.4", sz=17, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
tb(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.4),
   "437 节点  |  15 个 Frame 模块  |  20 个子群组  |  51 个参数", sz=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom cards
info = [("基准材质", "M_actor_pelica_cloth_04"), ("主节点群", "Arknights: Endfield_PBRToonBase"), ("代码产出", "hlsl/PBRToonBase.hlsl")]
for i, (label, val) in enumerate(info):
    x = Inches(1.8 + i * 3.5)
    add_rect(slide, x, Inches(5.5), Inches(3.1), Inches(1.0), BG_LIGHT, BORDER, Pt(1))
    tb(slide, x + Inches(0.15), Inches(5.55), Inches(2.8), Inches(0.3), label, sz=12, color=LIGHT_GRAY, bold=True)
    tb(slide, x + Inches(0.15), Inches(5.9), Inches(2.8), Inches(0.5), val, sz=12, color=DARK_GRAY, font="Consolas")


# ============================================================
# SLIDE 2: Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "1. 方案概述", "PBR + Toon 混合架构")

# Description
add_rect(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.3), BG_LIGHT, BORDER, Pt(1))
box = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.1))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PBRToonBase 是一套面向角色材质的混合渲染方案。以物理正确的微表面模型（GGX、能量守恒）为基础，通过 SigmoidSharp 阶跃函数和 Ramp 色带采样将光照响应风格化为卡通阴影，同时叠加边缘色、屏幕空间 Rim、视角色变等风格化效果。"
p.font.size = Pt(16)
p.font.color.rgb = MID_GRAY
p.font.name = "Microsoft YaHei"

# Three pillars
pillars = [
    ("PBR 物理基础", [
        "Cook-Torrance GGX 微表面模型",
        "Schlick Fresnel 菲涅尔",
        "Smith-Joint 可见性函数",
        "Kulla-Conty 能量守恒补偿",
        "FGD LUT 预积分查询",
    ]),
    ("Toon 风格化", [
        "SigmoidSharp 阶跃函数",
        "Ramp LUT 色带采样",
        "离散化阴影色带",
        "双路阴影（自阴影 + 投影）",
        "全局阴影亮度调节",
    ]),
    ("风格化叠加效果", [
        "Toon Fresnel 边缘色",
        "屏幕空间深度 Rim 光",
        "ThinFilm 视角色变（彩虹）",
        "SimpleTransmission 伪透射",
        "自发光 Emission",
    ]),
]
for i, (title, items) in enumerate(pillars):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(2.9), Inches(3.7), Inches(3.8), BG_LIGHT, BORDER, Pt(1))
    tb(slide, x + Inches(0.2), Inches(3.0), Inches(3.3), Inches(0.4), title, sz=20, color=BLACK, bold=True)
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), Inches(3.5), Inches(3.3), Inches(0.015))
    d.fill.solid()
    d.fill.fore_color.rgb = BG_MID
    d.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.3), Inches(3.6), Inches(3.2), Inches(3.0))
    tf2 = box.text_frame
    tf2.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = MID_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)

# Stats
stats = [("15", "Frame 模块"), ("20", "子群组"), ("51", "参数"), ("7", "功能开关")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.15)
    tb(slide, x, Inches(7.0), Inches(0.8), Inches(0.4), num, sz=26, color=BLACK, bold=True)
    tb(slide, x + Inches(0.85), Inches(7.05), Inches(2.0), Inches(0.35), label, sz=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 3: Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "2. 渲染管线总览", "15 个 Frame 模块按序组装")

# Left: main flow
stages = [
    ("贴图输入", "_D / _N / _P / _E 贴图解包"),
    ("初始化", "Frame.012 几何向量 N/V/L/T/B + 点积\nFrame.013 贴图拆通道 + PBR 参数"),
    ("前处理", "Frame.069 伪透射 — 修改 albedo"),
]
ys = [Inches(1.2), Inches(1.9), Inches(3.0)]
for (label, desc), y in zip(stages, ys):
    flow_box(slide, Inches(0.6), y, Inches(5.8), Inches(0.6) if '\n' not in desc else Inches(0.85), label + "\n" + desc, sz=12)

arrow_down(slide, Inches(3.5), Inches(1.8), Inches(0.15))
arrow_down(slide, Inches(3.5), Inches(2.75), Inches(0.3))

# Dual branch
flow_box(slide, Inches(0.6), Inches(4.1), Inches(2.7), Inches(1.0),
         "漫反射 DiffuseBRDF\nSigmoidSharp 阶跃\nRamp 色带采样", sz=12)
flow_box(slide, Inches(3.6), Inches(4.1), Inches(2.8), Inches(1.0),
         "高光反射 SpecularBRDF\nCook-Torrance GGX\n等向 / 各向异性可选", sz=12)

flow_box(slide, Inches(0.6), Inches(5.3), Inches(2.7), Inches(0.45),
         "ShadowAdjust 阴影调节", sz=12)
arrow_down(slide, Inches(1.9), Inches(5.1), Inches(0.2))

flow_box(slide, Inches(0.6), Inches(5.95), Inches(5.8), Inches(0.55),
         "IndirectLighting — FGD LUT + Kulla-Conty 能量守恒 + 直接光与间接光汇合", sz=12)
arrow_down(slide, Inches(1.9), Inches(5.75), Inches(0.2))
arrow_down(slide, Inches(5.0), Inches(5.1), Inches(0.85))

# Right: stylization effects
tb(slide, Inches(7.2), Inches(1.2), Inches(5.0), Inches(0.4),
   "风格化效果叠加（ADD）", sz=20, color=BLACK, bold=True)

effects = [
    ("Frame.008  ToonFresnel", "基于视角的 Toon 边缘色"),
    ("Frame.009  Rim", "屏幕空间深度边缘光，四因子遮罩"),
    ("Frame.010  Emission", "自发光通道"),
    ("Frame.011  ThinFilmFilter", "视角色变，RS LUT 彩虹效果"),
    ("Frame.014  Alpha", "最终透明度输出"),
]
for i, (frame, desc) in enumerate(effects):
    y = Inches(1.8 + i * 1.0)
    add_rect(slide, Inches(7.2), y, Inches(5.3), Inches(0.8), BG_LIGHT, BORDER, Pt(1))
    tb(slide, Inches(7.4), y + Inches(0.05), Inches(5.0), Inches(0.3), frame, sz=13, color=BLACK, bold=True, font="Consolas")
    tb(slide, Inches(7.4), y + Inches(0.38), Inches(5.0), Inches(0.35), desc, sz=12, color=MID_GRAY)
    if i < len(effects) - 1:
        arrow_down(slide, Inches(9.8), y + Inches(0.8), Inches(0.2))


# ============================================================
# SLIDE 4: Direct Diffuse
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "3.1 直接漫反射", "Frame.005 DiffuseBRDF + Frame.007 ShadowAdjust")

# Left: SigmoidSharp
tb(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
   "SigmoidSharp 阶跃函数", sz=20, color=BLACK, bold=True)

add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.3), BG_LIGHT, BORDER, Pt(1))
box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(2.1))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "将连续的 NdotL 光照响应转化为离散的阴影色带，实现手绘感卡通阴影。"
p.font.size = Pt(14)
p.font.color.rgb = MID_GRAY
p.font.name = "Microsoft YaHei"

pts = [
    "底数 = 100000（非自然指数），产生极其陡峭的 S 曲线",
    "center 控制阴影分界线位置（NdotL 到多大时开始过渡）",
    "sharp 控制边缘锐度（值越大，阴影边界越硬）",
    "与 smoothstep 不同，过渡区域可收窄到接近零",
]
for pt in pts:
    add_para(tf, pt, sz=13, color=MID_GRAY, sp=Pt(8))

# Right: dual shadow flow
tb(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.4),
   "双路阴影合并 + Ramp 采样", sz=20, color=BLACK, bold=True)

flow_box(slide, Inches(7.0), Inches(1.8), Inches(2.5), Inches(0.5), "自阴影 NoL", sz=12)
flow_box(slide, Inches(10.0), Inches(1.8), Inches(2.5), Inches(0.5), "投影阴影 CastShadow", sz=12)
arrow_down(slide, Inches(8.25), Inches(2.3), Inches(0.2))
arrow_down(slide, Inches(11.25), Inches(2.3), Inches(0.2))

flow_box(slide, Inches(7.0), Inches(2.6), Inches(2.5), Inches(0.45), "SigmoidSharp\n(center_HL, sharp_HL)", sz=11)
flow_box(slide, Inches(10.0), Inches(2.6), Inches(2.5), Inches(0.45), "SigmoidSharp\n(center_CS, sharp_CS)", sz=11)
arrow_down(slide, Inches(8.25), Inches(3.05), Inches(0.25))
arrow_down(slide, Inches(11.25), Inches(3.05), Inches(0.25))

flow_box(slide, Inches(8.2), Inches(3.4), Inches(3.8), Inches(0.45), "min() 取暗合并", sz=13)
arrow_down(slide, Inches(10.1), Inches(3.85), Inches(0.25))
flow_box(slide, Inches(7.5), Inches(4.2), Inches(5.0), Inches(0.45), "Ramp LUT 查表\nUV(X=阴影值, Y=0.5) -> RampSelect", sz=12)

# Ramp explanation card
bullet_card(slide, Inches(7.0), Inches(4.9), Inches(5.5), Inches(1.7), "Ramp 色带特点", [
    "Ramp LUT 支持 5 行，通过 _P.A 通道选择不同色带",
    "不同材质区域可使用不同 Ramp 行，一个 Shader 多种阴影风格",
    "rampAlpha 传递到 ShadowAdjust 驱动全局阴影亮度",
])

# Final formula
tb(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.35),
   "最终漫反射公式", sz=16, color=BLACK, bold=True)
formula_box(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.55),
            "directDiffuse = rampColor * diffuseColor * lightColor * directOcclusion", sz=14)

# directOcclusion note
add_rect(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.1), BG_LIGHT, BORDER, Pt(1))
box = slide.shapes.add_textbox(Inches(1.0), Inches(5.55), Inches(5.1), Inches(1.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "directOcclusion 由 AO 贴图和 _directOcclusionColor 混合得到，控制遮蔽区域的色调。投影阴影可独立调软硬，自阴影可以更柔和（半调渐变），投影保持硬边。"
p.font.size = Pt(13)
p.font.color.rgb = MID_GRAY
p.font.name = "Microsoft YaHei"


# ============================================================
# SLIDE 5: Specular
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "3.2 直接高光反射", "Frame.004 SpecularBRDF  |  Cook-Torrance F * D * V")

# Left: Isotropic
bullet_card(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(3.0), "等向 GGX（标准模式）", [
    "基于 Cook-Torrance 微表面模型：F * D * V",
    "F = Schlick Fresnel，fresnel0 由金属度在 0.04 和 albedo 间插值",
    "D = GGX 法线分布函数，单一粗糙度参数",
    "V = Smith-Joint 可见性函数",
    "适用于大多数非布料材质（金属、皮革等）",
    "各方向高光形状一致",
])

# Right: Anisotropic
bullet_card(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(3.0), "各向异性 GGX（布料/拉丝金属）", [
    "T/B 轴粗糙度分离，产生方向性拉长高光",
    "roughnessT / roughnessB 分别控制两个轴向",
    "适用于布料、拉丝金属等具有方向性纹理的材质",
    "沿纹理方向产生拉长的高光条纹",
    "SmoothnessMaxT / SmoothnessMaxB 独立可调",
])

# Three-level control
tb(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.4),
   "三级控制体系", sz=20, color=BLACK, bold=True)

levels = [
    ("1", "总开关 _UseAnisotropy", "关闭时直接使用等向 GGX"),
    ("2", "子开关 _UseToonAniso", "各向异性分支内的 Toon 变体选择"),
    ("3", "遮罩 _AnisotropicMask", "逐区域加权，材质不同部位使用不同高光模式"),
]
for i, (num, title, desc) in enumerate(levels):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(5.0), Inches(3.7), Inches(1.2), BG_LIGHT, BORDER, Pt(1))
    num_circle(slide, x + Inches(0.15), Inches(5.1), num)
    tb(slide, x + Inches(0.65), Inches(5.1), Inches(2.9), Inches(0.35), title, sz=14, color=BLACK, bold=True)
    tb(slide, x + Inches(0.2), Inches(5.55), Inches(3.3), Inches(0.55), desc, sz=13, color=MID_GRAY)

# Shared Fresnel note
add_rect(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.7), BG_LIGHT, BORDER, Pt(1))
box = slide.shapes.add_textbox(Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.6))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "两种高光模式共用同一 Schlick Fresnel：F_Schlick(fresnel0, 1, LdotH)，其中 fresnel0 = lerp(0.04, albedo, metallic)"
p.font.size = Pt(13)
p.font.color.rgb = MID_GRAY
p.font.name = "Microsoft YaHei"


# ============================================================
# SLIDE 6: Indirect Lighting
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "3.3 间接光照与能量守恒", "Frame.006 IndirectLighting")

# Three tasks
tasks = [
    ("1", "FGD LUT 预积分查询",
     "输入 (NdotV, perceptualRoughness) 查询 2D LUT",
     ["输出 specularFGD（预积分 Fresnel-GGX）",
      "输出 diffuseFGD（Disney 漫反射积分项）",
      "输出 reflectivity（总反射率 r）",
      "LUT 布局：R=B项, G=A+B(反射率), B=Disney Diffuse"]),
    ("2", "Kulla-Conty 能量守恒补偿",
     "标准微表面 BRDF 只考虑单次散射，粗糙表面多重散射能量损失可达 30%+",
     ["ecFactor = 1/reflectivity - 1",
      "光滑表面（r->1）：ecFactor->0，无需补偿",
      "粗糙表面（r~0.5）：ecFactor~1.0，补偿最强"]),
]

for i, (num, title, desc, items) in enumerate(tasks):
    x = Inches(0.8 + i * 6.15)
    num_circle(slide, x, Inches(1.2), num)
    tb(slide, x + Inches(0.5), Inches(1.2), Inches(5.0), Inches(0.35), title, sz=18, color=BLACK, bold=True)
    add_rect(slide, x, Inches(1.7), Inches(5.8), Inches(2.5), BG_LIGHT, BORDER, Pt(1))
    box = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.8), Inches(5.4), Inches(2.3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    for item in items:
        add_para(tf, item, sz=13, color=MID_GRAY, sp=Pt(6))

# Two compensation paths
paths = [
    ("A", "路径 A — 修正直接高光",
     "energyCompFactor = F0 * ecFactor + 1.0\n放大直接高光中因多重散射丢失的能量"),
    ("B", "路径 B — 补充间接高光",
     "indirectSpecComp = ecFactor * specularFGD\n以间接高光形式补充半球各方向散射能量"),
]
for i, (letter, title, desc) in enumerate(paths):
    x = Inches(0.8 + i * 6.15)
    num_circle(slide, x, Inches(4.5), letter)
    tb(slide, x + Inches(0.5), Inches(4.5), Inches(5.0), Inches(0.35), title, sz=16, color=BLACK, bold=True)
    add_rect(slide, x, Inches(5.0), Inches(5.8), Inches(1.0), BG_LIGHT, BORDER, Pt(1))
    box = slide.shapes.add_textbox(x + Inches(0.2), Inches(5.05), Inches(5.4), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = MID_GRAY
        p.font.name = "Microsoft YaHei" if j > 0 else "Consolas"
        p.space_before = Pt(4)

# Final
tb(slide, Inches(0.8), Inches(6.3), Inches(3.0), Inches(0.35), "最终光照汇合", sz=16, color=BLACK, bold=True)
formula_box(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
            "totalLighting = correctedSpecular + directDiffuse + indirectSpecComp + indirectDiffuse", sz=14)


# ============================================================
# SLIDE 7: Toon Fresnel + Rim
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "4.1-4.2 Toon Fresnel 边缘色 & 屏幕空间 Rim", "Frame.008 + Frame.009")

# Left: Toon Fresnel
bullet_card(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(3.0), "Toon Fresnel 边缘色", [
    "基于视角的边缘色调偏移，在角色轮廓处施加风格化色彩",
    "_ToonfresnelPow 控制 Fresnel 衰减速度",
    "SmoothStep(L, H) 裁切和软化边缘范围",
    "内/外双色混合：正面 fresnelInsideColor，掠射角过渡到 fresnelOutsideColor",
    "最终输出 = 颜色 * 权重，权重由视角驱动",
    "本质是对 PBR 结果做视角相关的色调偏移",
])

# Right: Rim
tb(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.4),
   "屏幕空间深度 Rim 光", sz=20, color=BLACK, bold=True)

tb(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
   "与传统 Fresnel Rim 不同，使用屏幕空间深度差检测作为主要轮廓判据", sz=13, color=MID_GRAY)

# Formula
formula_box(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(0.4),
            "rimMask = DepthRim * FresnelAtten * DirLightAtten * VerticalAtten", sz=13)

# Four factors
factors = [
    ("DepthRim", "法线方向偏移采样深度图\n深度突变处（物体边缘）产生 Rim"),
    ("FresnelAtten", "(1-NoV)^4\n掠射角加强，正对相机处抑制"),
    ("DirLightAtten", "lerp(atten, 1, saturate(NoL))\n光照方向响应，可控背光侧保留量"),
    ("VerticalAtten", "saturate(normalWS.y)\n顶部强、底部弱，防止脚底 Rim 过亮"),
]
for i, (name, desc) in enumerate(factors):
    x = Inches(7.0) if i % 2 == 0 else Inches(10.0)
    y = Inches(2.85) if i < 2 else Inches(4.35)
    add_rect(slide, x, y, Inches(2.7), Inches(1.3), BG_LIGHT, BORDER, Pt(1))
    tb(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.4), Inches(0.3),
       name, sz=14, color=BLACK, bold=True, font="Consolas")
    box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.4), Inches(2.4), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = MID_GRAY
        p.font.name = "Microsoft YaHei"

# Rim notes
bullet_card(slide, Inches(7.0), Inches(5.85), Inches(5.5), Inches(1.3), "Rim 特性", [
    "min(depthRim, 0.5) 截断贡献上限，防止过曝",
    "LoV 调制：光线背面加强、正面减弱，避免高光重叠",
    "最终 Rim = rimColor * rimMask",
])

# Left bottom: emission note
bullet_card(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.5), "自发光 Emission（Frame.010）", [
    "直接叠加 _E 贴图 RGB 到最终结果",
    "不受光照方向和阴影影响",
    "独立的自发光强度参数控制",
])


# ============================================================
# SLIDE 8: ThinFilm + Transmission
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "4.3-4.4 视角色变 & 伪透射", "Frame.011 ThinFilmFilter | Frame.069 SimpleTransmission")

# Left: ThinFilm
bullet_card(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(3.2), "ThinFilmFilter 视角色变", [
    "在角色边缘产生彩虹/光泽色变效果",
    "模拟薄膜干涉或特殊织物反光",
    "双 RS 贴图（Aurora / Yvonne）：不同角色使用不同色变纹理",
    "通过 _RS_Index 混合两张贴图",
    "Facing 驱动 1D LUT：仅用视角作为 X 轴（Y 固定 0.5）",
    "光照调制：阴影区域不叠加色变",
    "LIGHTEN 混合：max(base, rs) 取亮值，只增亮不减暗",
])

# Right: Transmission
bullet_card(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(3.2), "SimpleTransmission 伪透射", [
    "通过屏幕空间采样模拟半透明材质的透射效果",
    "在管线最前端执行，修改 albedo 后影响后续所有漫反射",
    "菲涅尔偏移：边缘处采样到更远的背景，产生折射扭曲",
    "IOR = 1.25，默认混合强度 0.65",
    "采样 _CameraOpaqueTexture 获取背景场景色",
    "lerp(sceneColor, albedo, _SimpleTransmissionValue) 混合",
    "_UseSimpleTransmission 开关控制",
])

# Process flow at bottom
tb(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.35),
   "伪透射在管线中的位置", sz=16, color=BLACK, bold=True)

boxes = [
    "SimpleTransmission\n修改 albedo",
    "DiffuseBRDF\nSigmoidSharp + Ramp",
    "SpecularBRDF\nGGX 高光",
    "IndirectLighting\n能量守恒汇合",
    "风格化叠加\nFresnel/Rim/ThinFilm",
]
for i, text in enumerate(boxes):
    x = Inches(0.6 + i * 2.5)
    bg = BG_MID if i == 0 else BG_LIGHT
    flow_box(slide, x, Inches(5.2), Inches(2.2), Inches(0.85), text, fill=bg, sz=11)
    if i < len(boxes) - 1:
        arrow_right(slide, x + Inches(2.2), Inches(5.62), Inches(0.3))

tb(slide, Inches(0.6), Inches(6.2), Inches(11.7), Inches(0.35),
   "albedo 在管线入口被修改 -> 后续所有漫反射计算都基于混合后的 albedo", sz=13, color=MID_GRAY)


# ============================================================
# SLIDE 9: Design Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "5. 设计特点", "PBR-Toon 平衡 | 高度参数化 | 模块化架构")

features = [
    ("PBR 与 Toon 的平衡", [
        ("高光反射", "保留完整 Cook-Torrance 微表面模型"),
        ("间接光照", "保留 FGD 预积分和能量守恒"),
        ("漫反射", "通过 SigmoidSharp + Ramp 风格化"),
        ("效果", "PBR 光照环境中物理一致 + 卡通手绘感"),
    ]),
    ("高度参数化", [
        ("51 个参数", "覆盖材质各个维度"),
        ("7 个开关", "按需启用/禁用特定效果"),
        ("逐区域控制", "遮罩加权，不同部位不同效果"),
        ("美术友好", "直觉化的面向美术参数设计"),
    ]),
    ("模块化 Frame 架构", [
        ("15 个 Frame", "各自独立，通过结构体传递数据"),
        ("独立调试", "单个模块可独立调试和替换"),
        ("可扩展", "新效果作为新 Frame 插入，不影响已有模块"),
        ("20 个子群组", "SigmoidSharp、F_Schlick 等跨模块复用"),
    ]),
]

for col, (title, items) in enumerate(features):
    x = Inches(0.6 + col * 4.2)
    tb(slide, x, Inches(1.2), Inches(3.8), Inches(0.4), title, sz=20, color=BLACK, bold=True)
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.7), Inches(3.8), Inches(0.015))
    d.fill.solid()
    d.fill.fore_color.rgb = BLACK
    d.line.fill.background()
    for i, (label, desc) in enumerate(items):
        y = Inches(1.9 + i * 1.15)
        add_rect(slide, x, y, Inches(3.8), Inches(1.0), BG_LIGHT, BORDER, Pt(1))
        tb(slide, x + Inches(0.15), y + Inches(0.08), Inches(3.5), Inches(0.3),
           label, sz=14, color=BLACK, bold=True)
        tb(slide, x + Inches(0.15), y + Inches(0.45), Inches(3.5), Inches(0.5),
           desc, sz=13, color=MID_GRAY)

# Switch table
tb(slide, Inches(0.6), Inches(6.5), Inches(1.5), Inches(0.35),
   "功能开关：", sz=14, color=BLACK, bold=True)
switches = [
    ("UseAnisotropy", "各向异性高光"),
    ("UseToonAniso", "Toon 变体"),
    ("UseNormalTex", "法线贴图"),
    ("UseRSEff", "视角色变"),
    ("UseSimpleTransmission", "伪透射"),
    ("IsSkin", "皮肤材质"),
]
for i, (sw, desc) in enumerate(switches):
    x = Inches(2.1 + i * 1.9)
    add_rect(slide, x, Inches(6.45), Inches(1.75), Inches(0.7), BG_LIGHT, BORDER, Pt(1))
    tb(slide, x + Inches(0.1), Inches(6.47), Inches(1.55), Inches(0.3), sw, sz=10, color=BLACK, bold=True, font="Consolas")
    tb(slide, x + Inches(0.1), Inches(6.78), Inches(1.55), Inches(0.3), desc, sz=10, color=MID_GRAY)


# ============================================================
# SLIDE 10: Texture Conventions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "6. 贴图通道约定", "PBRToonBase 贴图打包规范")

textures = [
    ("_D", "RGB", "Albedo 颜色", "sRGB"),
    ("_D", "A", "Alpha 遮罩", "—"),
    ("_N", "RG", "切线空间法线 XY", "Non-Color"),
    ("_P", "R", "Metallic 金属度", "Non-Color"),
    ("_P", "G", "AO / directOcclusion", "Non-Color"),
    ("_P", "B", "Smoothness 感知平滑度", "Non-Color"),
    ("_P", "A", "RampUV（Toon Ramp 行选择）", "Non-Color"),
    ("_E", "RGB", "Emission 自发光", "Non-Color"),
    ("_M", "RGBA", "遮罩（部分材质使用）", "Non-Color"),
    ("RS", "—", "视角色变 LUT（Aurora / Yvonne）", "Linear/sRGB"),
]

# Header
hy = Inches(1.3)
cw = [Inches(1.5), Inches(1.2), Inches(5.0), Inches(2.0)]
cx = [Inches(1.5)]
for w in cw[:-1]:
    cx.append(cx[-1] + w)

hbar = add_rect(slide, Inches(1.3), hy, Inches(9.7), Inches(0.5), BLACK, rounded=False)
for header, x, w in zip(["贴图", "通道", "语义", "色彩空间"], cx, cw):
    tb(slide, x, hy + Inches(0.07), w, Inches(0.35), header, sz=14, color=BG_WHITE, bold=True)

for ri, (tex, ch, sem, cs) in enumerate(textures):
    y = Inches(1.85 + ri * 0.48)
    bg = BG_LIGHT if ri % 2 == 0 else BG_WHITE
    add_rect(slide, Inches(1.3), y, Inches(9.7), Inches(0.46), bg, rounded=False)
    vals = [tex, ch, sem, cs]
    for ci, (val, x, w) in enumerate(zip(vals, cx, cw)):
        fn = "Consolas" if ci < 2 else "Microsoft YaHei"
        c = BLACK if ci == 0 else DARK_GRAY
        b = (ci == 0)
        tb(slide, x, y + Inches(0.05), w, Inches(0.35), val, sz=12, color=c, bold=b, font=fn)

# Note
tb(slide, Inches(1.5), Inches(6.8), Inches(9.0), Inches(0.35),
   "_P 贴图：4 通道 = Metallic + AO + Smoothness + RampUV（最大化通道打包利用率）", sz=14, color=BLACK, bold=True)


# ============================================================
# SLIDE 11: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLACK
bar.line.fill.background()

tb(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.9),
   "Thank You", sz=56, color=BLACK, bold=True, font="Consolas", align=PP_ALIGN.CENTER)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(2.333), Inches(0.025))
div.fill.solid()
div.fill.fore_color.rgb = BLACK
div.line.fill.background()

tb(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
   "PBRToonBase  |  PBR + Toon 混合渲染管线", sz=20, color=DARK_GRAY, align=PP_ALIGN.CENTER)
tb(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5),
   "Arknights: Endfield  /  Goo Engine 4.4", sz=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Summary cards
summary = [
    ("物理基础", "Cook-Torrance GGX\n能量守恒补偿"),
    ("Toon 风格", "SigmoidSharp 阶跃\nRamp 色带采样"),
    ("风格化效果", "Fresnel 边缘色 + Rim\nThinFilm + 伪透射"),
    ("模块化架构", "15 个 Frame 模块\n20 个复用子群组"),
]
for i, (title, desc) in enumerate(summary):
    x = Inches(1.5 + i * 2.8)
    add_rect(slide, x, Inches(5.2), Inches(2.5), Inches(1.3), BG_LIGHT, BORDER, Pt(1))
    tb(slide, x + Inches(0.15), Inches(5.3), Inches(2.2), Inches(0.35), title, sz=14, color=BLACK, bold=True)
    box = slide.shapes.add_textbox(x + Inches(0.15), Inches(5.7), Inches(2.2), Inches(0.7))
    tf2 = box.text_frame
    tf2.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = MID_GRAY
        p.font.name = "Microsoft YaHei"


# Save
output_path = r"f:\GooBlenderProj\shader-migration\PBRToonBase_TechSharing.pptx"
prs.save(output_path)
print("PPT saved to: " + output_path)
